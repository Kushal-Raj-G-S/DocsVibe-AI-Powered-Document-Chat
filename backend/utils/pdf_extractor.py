"""
Document Text Extraction
Supports PDF, DOCX, and PPTX file formats
Extracts directly from bytes (no temp files)
"""

import PyPDF2
from typing import Tuple, Optional, Union
import logging
import io

logger = logging.getLogger(__name__)

# Optional imports for DOCX and PPTX
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available, DOCX support disabled")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available, PPTX support disabled")


def extract_text_from_pdf(pdf_data: Union[str, bytes]) -> Tuple[str, int]:
    """
    Extract text from PDF file or bytes
    
    Args:
        pdf_data: Path to PDF file or PDF file bytes
        
    Returns:
        Tuple of (extracted_text, page_count)
    """
    try:
        if isinstance(pdf_data, bytes):
            # Extract from bytes (in-memory)
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_data))
        else:
            # Extract from file path
            with open(pdf_data, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
        
        page_count = len(pdf_reader.pages)
        
        text = ""
        for page_num in range(page_count):
            page = pdf_reader.pages[page_num]
            text += f"\n--- Page {page_num + 1} ---\n"
            text += page.extract_text()
        
        logger.info(f"✅ Extracted text from PDF: {page_count} pages")
        return text.strip(), page_count
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        raise Exception(f"Error extracting PDF text: {str(e)}")


def extract_text_from_docx(docx_data: Union[str, bytes]) -> Tuple[str, int]:
    """
    Extract text from DOCX file or bytes
    
    Args:
        docx_data: Path to DOCX file or DOCX file bytes
        
    Returns:
        Tuple of (extracted_text, paragraph_count)
    """
    if not DOCX_AVAILABLE:
        raise Exception("DOCX support not available. Install python-docx package.")
    
    try:
        if isinstance(docx_data, bytes):
            # Extract from bytes (in-memory)
            doc = Document(io.BytesIO(docx_data))
        else:
            # Extract from file path
            doc = Document(docx_data)
        
        text_parts = []
        paragraph_count = 0
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
                paragraph_count += 1
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        text = "\n\n".join(text_parts)
        logger.info(f"✅ Extracted text from DOCX: {paragraph_count} paragraphs, {len(doc.tables)} tables")
        return text.strip(), paragraph_count
        
    except Exception as e:
        logger.error(f"Error extracting DOCX text: {e}")
        raise Exception(f"Error extracting DOCX text: {str(e)}")


def extract_text_from_pptx(pptx_data: Union[str, bytes]) -> Tuple[str, int]:
    """
    Extract text from PPTX file or bytes
    
    Args:
        pptx_data: Path to PPTX file or PPTX file bytes
        
    Returns:
        Tuple of (extracted_text, slide_count)
    """
    if not PPTX_AVAILABLE:
        raise Exception("PPTX support not available. Install python-pptx package.")
    
    try:
        if isinstance(pptx_data, bytes):
            # Extract from bytes (in-memory)
            prs = Presentation(io.BytesIO(pptx_data))
        else:
            # Extract from file path
            prs = Presentation(pptx_data)
            
        slide_count = len(prs.slides)
        
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            
            if slide_text.strip() != f"--- Slide {slide_num} ---":
                text_parts.append(slide_text)
        
        text = "\n".join(text_parts)
        logger.info(f"✅ Extracted text from PPTX: {slide_count} slides")
        return text.strip(), slide_count
        
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {e}")
        raise Exception(f"Error extracting PPTX text: {str(e)}")


def extract_text_from_file(file_data: Union[str, bytes], file_type: Optional[str] = None) -> Tuple[str, int, str]:
    """
    Extract text from any supported document type
    Supports both file paths and bytes (in-memory)
    
    Args:
        file_data: Path to file or file bytes
        file_type: File type hint ('pdf', 'docx', 'pptx')
        
    Returns:
        Tuple of (extracted_text, unit_count, detected_type)
    """
    if file_type == 'pdf':
        text, count = extract_text_from_pdf(file_data)
        return text, count, 'pdf'
    elif file_type == 'docx':
        text, count = extract_text_from_docx(file_data)
        return text, count, 'docx'
    elif file_type == 'pptx':
        text, count = extract_text_from_pptx(file_data)
        return text, count, 'pptx'
    else:
        raise Exception(f"Unsupported file type: {file_type}")

